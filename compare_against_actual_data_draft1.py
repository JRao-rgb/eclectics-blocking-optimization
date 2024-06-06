# -*- coding: utf-8 -*-
"""
Created on Sun Jun  2 19:20:56 2024

@author: jraos
"""

# -*- coding: utf-8 -*-
"""
Created on Sun Apr 14 10:26:03 2024

@author: jraos
"""

import pptx as ppt
import os
import numpy as np
import matplotlib.pyplot as plt
os.chdir("C:\\Users\\jraos\\OneDrive - Stanford\\Documents\\Stanford\\EC\\Eclectics Blocking Optimization\\eclectics-blocking-optimization")
import computational_geometry as geo
import optimize
import metrics

os.chdir("C:/Users/jraos/OneDrive - Stanford/Documents/Stanford/EC/Eclectics Blocking Optimization")

data_dir = "data"
slide_width = 13.3333 # the slide width, in inches. Used to normnalize the dancer positions
slide_height = 7.5 # the slide hieght, in inches. Used to normalized the dancer positions
blocking_diagram = ppt.Presentation(data_dir + "/wake-up-cleaned.pptx")

#%% extracting all the text from the slides

text_runs = []
text_positions = []

for slide in blocking_diagram.slides:
    text_positions_per_slide = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
                text_positions_per_slide.append([run.text, shape.left.inches, shape.top.inches])
    text_positions.append(text_positions_per_slide)

#%% evaluate metrics for the human-made formation 

num_dancers = len(text_positions[0])
num_formations = len(text_positions)

X = np.zeros((num_dancers,num_formations))
Y = np.zeros((num_dancers,num_formations))
P_initial = np.zeros((num_dancers,num_formations),dtype = np.int8)
P_inv = np.zeros((num_dancers,num_formations), dtype = np.int8)

# first, create a mapping from dancerID to positionID
M_inv = np.zeros(100)

for positionID in range(num_dancers):
    dancerID = int(text_positions[0][positionID][0])
    M_inv[dancerID] = positionID

# first, convert it into the familiar for of X, Y, P
for formation_index, formation in enumerate(text_positions):
    for coordinateID, dancer in enumerate(formation):
        X[coordinateID, formation_index] = dancer[1]
        Y[coordinateID, formation_index] = dancer[2]
        P_initial[coordinateID, formation_index] = M_inv[int(dancer[0])]
        P_inv[P_initial[coordinateID, formation_index], formation_index] = coordinateID

# rescale X, Y
X = (X - slide_width/2) / slide_width * 20
Y = (Y - slide_height/2) / slide_height * 10

#%% generate animations for the human-made formation

os.chdir("C:\\Users\\jraos\\OneDrive - Stanford\\Documents\\Stanford\\EC\\Eclectics Blocking Optimization\\eclectics-blocking-optimization")
geo.animate_movement(X, Y, P_inv, filename = "6-2-24_animations/wake-up-cleaned-human.gif")

#%% compute metrics for this formation

print("number of intersections:",metrics.num_intersections(X, Y, P_inv))
print("average distance:",metrics.average_distance(X, Y, P_inv))
print("maximum distance:",metrics.max_distance(X, Y, P_inv))

#%% see if my code can do any better

sf1 = 1
sf2 = 1
sf3 = 1
max_iteration = 10
C = np.full((num_dancers, num_formations),-1,dtype=np.int8)
P_init = np.full((num_dancers, num_formations), 0, dtype=np.int8)
P_init[:,] = np.linspace(0,num_dancers-1,num_dancers)[...,None]
P, P_inv, M, M_inv, cost_array = optimize.optimize_formation(X, Y, C, P_init, max_iteration, sf1, sf2, sf3)

print("number of intersections:",metrics.num_intersections(X, Y, P_inv))
print("average distance:",metrics.average_distance(X, Y, P_inv))
print("maximum distance:",metrics.max_distance(X, Y, P_inv))

#%%

os.chdir("C:\\Users\\jraos\\OneDrive - Stanford\\Documents\\Stanford\\EC\\Eclectics Blocking Optimization\\eclectics-blocking-optimization")
geo.animate_movement(X, Y, P_inv, filename = "6-2-24_animations/wake-up-computer.gif")
