# -*- coding: utf-8 -*-
"""
Created on Sun Apr 14 10:26:03 2024

@author: jraos
"""

import pptx as ppt
import os

os.chdir("C:/Users/jraos/OneDrive - Stanford/Documents/Stanford/EC/Eclectics Blocking Optimization")

data_dir = "data"
slide_width = 13.3333 # the slide width, in inches. Used to normnalize the dancer positions
slide_height = 7.5 # the slide hieght, in inches. Used to normalized the dancer positions
blocking_diagram = ppt.Presentation(data_dir + "/wake-up-cleaned.pptx")

#%% extracting all the text from the slides

text_runs = []
text_positions = []

for slide in blocking_diagram.slides:
    text_positions_per_slide = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                text_runs.append(run.text)
                text_positions_per_slide.append([run.text, shape.left.inches, shape.top.inches])
    text_positions.append(text_positions_per_slide)

#%% displaying the text at the correct locations

import matplotlib.pyplot as plt

plt.ion()
plt.figure()
for formation in text_positions:
    for dancer in formation:
        plt.text(dancer[1]/slide_width,dancer[2]/slide_height,dancer[0],size='large')
    break

#%% animating the formations as each dancer moves from one to another

import numpy as np

position_matrix = -1*np.ones((60,20,2)) # array that contains the dancerID in the rows and the 
# timed position in the columns. So #rows = #dancers and #cols = #formations
# dancer0 is a test object because Python is 0-indexed

for i, formation in enumerate(text_positions):
    for j, dancer in enumerate(formation):
        position_matrix[int(dancer[0])][i][0] = dancer[1]/slide_width
        position_matrix[int(dancer[0])][i][1] = dancer[2]/slide_height

# now, make it so that each dancer's position is padded and it kind of animates
# the movement of the dancers across the stage

fps = 10 # 10 frames between each formation in the animation
freeze_frames = 10 # number of frames the current positions will be frozen for
animated_position_matrix = -1 * np.ones((60,(20-1)*fps+20*freeze_frames,2))

for i in range(np.shape(position_matrix)[0]-1):
    # create the initial freeze frames
    for k in range(freeze_frames):
        animated_position_matrix[i][k][0] = position_matrix[i,0,0]
        animated_position_matrix[i][k][1] = position_matrix[i,0,1]
        
    for j in range(np.shape(position_matrix)[1]-1):
        x1 = position_matrix[i,j,0]
        x2 = position_matrix[i,j+1,0]
        y1 = position_matrix[i,j,1]
        y2 = position_matrix[i,j+1,1]
        
        if int(y1) == -1:
            x1 = int(x2 > 0.5)
            y1 = y2
        elif int(y2) == -1:
            x2 = int(x1 > 0.5)
            y2 = y1
        
        for k in range(fps):
            animated_position_matrix[i][j*(fps+freeze_frames)+freeze_frames+k][0] = x1 + k * (x2-x1)/fps
            animated_position_matrix[i][j*(fps+freeze_frames)+freeze_frames+k][1] = y1 + k * (y2-y1)/fps
        
        for k in range(freeze_frames):
            animated_position_matrix[i][j*(fps+freeze_frames)+fps+freeze_frames+k][0] = x2
            animated_position_matrix[i][j*(fps+freeze_frames)+fps+freeze_frames+k][1] = y2
            
#%% making a video of the whole thing

import matplotlib.animation as animation

# Create an array of images (for demonstration purposes)
# You can replace this with your actual images
num_frames = (20-1)*fps+1

# Create a figure
plt.ioff()
fig, ax = plt.subplots()

# Initialize an empty list to store frames
frames = []

# Add images to frames
for i in range(num_frames):
    container = ax.scatter(animated_position_matrix[:,i,0], animated_position_matrix[:,i,1], animated=True, color='b')
    plt.xlim(0,1)
    plt.ylim(0,1)
    frames.append([container])

# Create an animation
ani = animation.ArtistAnimation(fig=fig, artists=frames, interval=100, blit=True, repeat_delay=1000)

# Save the animation as a video
ani.save("output_video.gif", writer="pillow")

#%% extracting all the shapes fromm the slides

shape_list = []

for slide in blocking_diagram.slides:
    shape_list.append(slide.shapes)
    for shape in slide.shapes:
        print(shape.top, shape.left, shape.height, shape.width)